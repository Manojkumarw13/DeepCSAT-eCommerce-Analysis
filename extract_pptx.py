from pptx import Presentation

try:
    prs = Presentation('p:/Project/Labmentix/Deep_Cast/DeepCSAT – Ecommerce.pptx')
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    with open('p:/Project/Labmentix/Deep_Cast/pptx_content.txt', 'w', encoding='utf-8') as f:
        f.write('\n\n'.join(text))
    print("Successfully extracted pptx content to pptx_content.txt")
except Exception as e:
    print(f"Error: {e}")
